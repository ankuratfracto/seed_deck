from pathlib import Path
import re
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "road-transport-sales-deck-editable-dynamic.pptx"


COLORS = {
    "ink": RGBColor(20, 33, 61),
    "ink_soft": RGBColor(50, 68, 99),
    "paper": RGBColor(255, 253, 248),
    "panel": RGBColor(255, 255, 255),
    "panel_soft": RGBColor(245, 247, 251),
    "line": RGBColor(215, 222, 234),
    "accent": RGBColor(14, 143, 125),
    "accent_dark": RGBColor(11, 109, 96),
    "accent_soft": RGBColor(218, 246, 240),
    "signal": RGBColor(242, 163, 59),
    "signal_soft": RGBColor(255, 241, 219),
    "berry": RGBColor(191, 91, 118),
    "navy": RGBColor(21, 49, 77),
    "white": RGBColor(255, 255, 255),
}


prs = Presentation()
prs.slide_width = Inches(11.69)
prs.slide_height = Inches(8.27)
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["paper"]


def set_run(run, size=12, bold=False, color="ink", font="IBM Plex Sans"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]


def add_text(slide, x, y, w, h, text, size=12, bold=False, color="ink", align=PP_ALIGN.LEFT, font="IBM Plex Sans"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, color=color, font=font)
    return shape


def add_multiline(slide, x, y, w, h, lines, size=12, color="ink_soft", bullet=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {line}" if bullet else line
        p.space_after = Pt(5)
        for run in p.runs:
            set_run(run, size=size, color=color)
    return shape


def add_panel(slide, x, y, w, h, fill="panel", line="line", radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    shape.line.width = Pt(1)
    shape.adjustments[0] = radius
    return shape


def add_band(slide, x, y, w, h, title, body):
    shape = add_panel(slide, x, y, w, h, fill="navy", line="navy")
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.35, title, size=18, bold=True, color="white", font="Space Grotesk")
    add_text(slide, x + 0.18, y + 0.58, w - 0.36, h - 0.7, body, size=11.2, color="white")
    return shape


def add_tag(slide, x, y, text, w=1.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["signal_soft"]
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    add_text(slide, x + 0.05, y + 0.055, w - 0.1, 0.18, text, size=7.8, bold=True, color="accent_dark", align=PP_ALIGN.CENTER)
    return shape


def add_eyebrow(slide, x, y, text, w=2.6):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["accent_soft"]
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    add_text(slide, x + 0.12, y + 0.07, w - 0.24, 0.2, text, size=8.4, bold=True, color="accent_dark", align=PP_ALIGN.CENTER)
    return shape


def add_chip(slide, x, y, text, w):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["panel"]
    shape.line.color.rgb = COLORS["line"]
    shape.adjustments[0] = 0.5
    add_text(slide, x + 0.05, y + 0.075, w - 0.1, 0.18, text, size=8.5, bold=True, color="ink_soft", align=PP_ALIGN.CENTER)
    return shape


def add_card(slide, x, y, w, h, tag, title, body, tag_w=0.72):
    add_panel(slide, x, y, w, h)
    add_tag(slide, x + 0.18, y + 0.16, tag, w=tag_w)
    add_text(slide, x + 0.18, y + 0.58, w - 0.36, 0.42, title, size=12.2, bold=True, color="ink")
    add_text(slide, x + 0.18, y + 1.06, w - 0.36, h - 1.14, body, size=8.6, color="ink_soft")


def add_metric(slide, x, y, w, h, value, label):
    add_panel(slide, x, y, w, h, fill="panel")
    add_text(slide, x + 0.15, y + 0.16, w - 0.3, 0.54, value, size=11.6, bold=True, color="accent_dark")
    add_text(slide, x + 0.15, y + 0.76, w - 0.3, h - 0.88, label, size=7.9, color="ink_soft")


def add_header(slide, section):
    if (ROOT / "fracto_logo.png").exists():
        slide.shapes.add_picture(str(ROOT / "fracto_logo.png"), Inches(0.55), Inches(0.28), width=Inches(0.72), height=Inches(0.225))
    add_text(slide, 9.1, 0.36, 2.0, 0.25, section, size=8.3, bold=True, color="ink_soft", align=PP_ALIGN.RIGHT)


def add_footer(slide, page):
    add_text(slide, 0.55, 7.78, 3.2, 0.24, "Prepared for {{lead_company_name}}", size=8.5, color="ink_soft")
    add_text(slide, 10.55, 7.78, 0.55, 0.24, f"{page:02d}", size=8.5, color="ink_soft", align=PP_ALIGN.RIGHT)


def new_slide(section, page):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_header(slide, section)
    add_footer(slide, page)
    return slide


def add_logo_card(slide, x, y, w, h, img_path, fallback):
    add_panel(slide, x, y, w, h)
    path = ROOT / img_path
    if path.exists():
        with Image.open(path) as img:
            img_w, img_h = img.size
        pad_x, pad_y = 0.12, 0.1
        box_w, box_h = w - (pad_x * 2), h - (pad_y * 2)
        img_ratio = img_w / img_h
        box_ratio = box_w / box_h
        if img_ratio > box_ratio:
            out_w = box_w
            out_h = box_w / img_ratio
        else:
            out_h = box_h
            out_w = box_h * img_ratio
        out_x = x + (w - out_w) / 2
        out_y = y + (h - out_h) / 2
        slide.shapes.add_picture(str(path), Inches(out_x), Inches(out_y), width=Inches(out_w), height=Inches(out_h))
    else:
        add_text(slide, x + 0.15, y + 0.37, w - 0.3, 0.32, fallback, size=11, bold=True, color="ink_soft", align=PP_ALIGN.CENTER)


def add_image_contain(slide, img_path, x, y, w, h):
    path = ROOT / img_path
    with Image.open(path) as img:
        img_w, img_h = img.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        out_w = w
        out_h = w / img_ratio
    else:
        out_h = h
        out_w = h * img_ratio
    out_x = x + (w - out_w) / 2
    out_y = y + (h - out_h) / 2
    return slide.shapes.add_picture(str(path), Inches(out_x), Inches(out_y), width=Inches(out_w), height=Inches(out_h))


def build_slide_1():
    slide = new_slide("Road Transportation Deck", 1)
    add_eyebrow(slide, 0.72, 1.02, "AI INFRASTRUCTURE FOR TRANSPORT DOCUMENT OPS", w=3.95)
    add_text(slide, 0.72, 1.58, 5.95, 1.28, "Automate {{primary_workflow}} for {{lead_company_name}}.", size=25, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.92, 5.95, 0.72, "Prepared for {{prospect_name}} and the {{lead_company_name}} team. Fracto can classify transport documents, extract key fields, and check whether signatures, stamps, and handwritten remarks are present.", size=11.4, color="ink_soft")
    add_panel(slide, 0.72, 3.88, 5.45, 1.22)
    add_multiline(slide, 0.93, 4.07, 4.95, 0.82, [
        "{{sample_document_types}}",
        "Signature and stamp presence checks",
        "Handwritten remark detection",
        "APIs for ops and proof workflows"
    ], size=9.8, bullet=True)
    for x, y, text, width in [
        (0.72, 5.35, "PODs", 0.78),
        (2.67, 5.35, "Trip sheet review", 1.45),
        (4.27, 5.35, "Hire letters", 1.1),
        (0.72, 5.78, "Remarks", 0.9),
        (1.78, 5.78, "DL, RC, Aadhaar, PAN", 1.95),
    ]:
        add_chip(slide, x, y, text, width)
    add_panel(slide, 7.0, 1.02, 3.95, 4.72, fill="panel_soft")
    add_text(slide, 7.27, 1.26, 3.25, 0.58, "What this changes for {{lead_company_name}}", size=14.5, bold=True, color="ink")
    add_metric(slide, 7.27, 2.02, 1.55, 1.32, "POD ready", "Works on proof-of-delivery and transport movement documents")
    add_metric(slide, 9.05, 2.02, 1.55, 1.32, "Presence checks", "Detect signature, stamp, and handwritten notes")
    add_metric(slide, 7.27, 3.5, 1.55, 1.32, "Ops ready", "Structured outputs plus review signals")
    add_metric(slide, 9.05, 3.5, 1.55, 1.32, "Faster proofing", "Find incomplete documents before downstream disputes")
    add_band(slide, 7.05, 5.92, 3.9, 1.12, "Positioning", "Replace manual proof checking with one document layer for {{lead_company_name}}.")


def build_slide_2():
    slide = new_slide("Why Transport Ops Break Down", 3)
    add_eyebrow(slide, 0.72, 1.0, "THE OPERATIONAL PROBLEM", w=2.35)
    add_text(slide, 0.72, 1.48, 6.0, 0.98, "Transport document review is still manual, repetitive, and easy to miss.", size=21.5, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.58, 5.9, 0.48, "For {{lead_company_name}}, {{pain_point_1}} can slow down {{primary_workflow}}.", size=10.9, color="ink_soft")
    cards = [
        ("POD", "Signoff checks are manual", "Teams still inspect whether a POD is signed, stamped, or annotated."),
        ("Exceptions", "Handwritten notes get missed", "Delivery remarks and shortages often sit in handwritten sections."),
        ("Finance", "Trip sheets and hire docs need structure", "Core fields still get keyed manually for matching and review."),
        ("Scale", "Proof queues pile up fast", "Backlogs build when every document needs a human visual check."),
    ]
    coords = [(0.72, 3.28), (3.55, 3.28), (0.72, 5.1), (3.55, 5.1)]
    for (tag, title, body), (x, y) in zip(cards, coords):
        add_card(slide, x, y, 2.55, 1.58, tag, title, body, tag_w=1.0)
    add_panel(slide, 6.85, 1.52, 3.95, 2.2)
    add_text(slide, 7.08, 1.78, 3.42, 0.33, "Typical failure pattern", size=17, bold=True, color="ink")
    add_multiline(slide, 7.1, 2.24, 3.35, 1.05, [
        "PODs and trip sheets arrive in mixed quality",
        "Teams check stamps and signatures by eye",
        "Handwritten remarks are missed or delayed",
        "Proof completion slows billing and follow-up",
    ], size=10.3, bullet=True)
    add_panel(slide, 6.85, 4.1, 3.95, 2.12, fill="panel_soft")
    add_text(slide, 7.08, 4.35, 3.42, 0.33, "What operators need instead", size=17, bold=True, color="ink")
    add_multiline(slide, 7.1, 4.8, 3.35, 1.0, [
        "Fast parsing plus visual presence checks",
        "Flags for missing signoff or missing stamp",
        "Handwritten remark detection for exceptions",
        "Review only the uncertain cases",
    ], size=10.3, bullet=True)


def build_slide_3():
    slide = new_slide("Capabilities", 3)
    add_eyebrow(slide, 0.72, 1.0, "CORE MODULES", w=1.65)
    add_text(slide, 0.72, 1.48, 8.7, 0.72, "Capabilities mapped to {{lead_company_name}} proof and ops workflows.", size=21.5, bold=True, color="ink", font="Space Grotesk")
    cards = [
        ("01", "{{use_case_1_title}}", "{{use_case_1_copy}}"),
        ("02", "{{use_case_2_title}}", "{{use_case_2_copy}}"),
        ("03", "{{use_case_3_title}}", "{{use_case_3_copy}}"),
        ("04", "Signature detection", "Check whether a required signature is present."),
        ("05", "Stamp detection", "Flag whether the proof carries a visible stamp."),
        ("06", "Invoice parser", "Pull invoice numbers, dates, and amount fields when needed downstream."),
    ]
    x0, y0 = 0.72, 2.35
    for idx, (tag, title, body) in enumerate(cards):
        x = x0 + (idx % 3) * 3.45
        y = y0 + (idx // 3) * 1.82
        add_card(slide, x, y, 3.1, 1.52, tag, title, body, tag_w=0.58)
    add_panel(slide, 0.72, 6.12, 10.05, 0.95, fill="panel_soft")
    add_text(slide, 0.95, 6.3, 1.8, 0.24, "Workflow bundle", size=12, bold=True, color="ink")
    cx, cy = 2.55, 6.27
    for text, width in [
        ("Upload", 0.72), ("Classify", 0.82), ("Extract", 0.72), ("Check signature", 1.24),
        ("Check stamp", 1.06), ("Flag remarks", 1.1), ("Review", 0.76), ("Approve", 0.8)
    ]:
        if cx + width > 10.45:
            cx, cy = 2.55, cy + 0.38
        add_chip(slide, cx, cy, text, width)
        cx += width + 0.08


def add_signal_block(slide, x, y, w, h, label, value):
    add_panel(slide, x, y, w, h, fill="panel")
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.22, label, size=8.2, bold=True, color="accent_dark")
    add_text(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.58, value, size=9.2, color="ink_soft")


def build_slide_4_visual_proof():
    slide = new_slide("Visual Proof", 4)
    add_eyebrow(slide, 0.72, 1.0, "DOCUMENT IMAGE + EXTRACTION", w=2.75)
    add_text(slide, 0.72, 1.48, 8.6, 0.82, "Visual proof: raw POD image to structured extraction and completeness checks.", size=20.5, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.28, 8.2, 0.34, "Representative transport proof. Fracto keeps the document visible while returning fields, confidence, and visual signoff signals.", size=10.4, color="ink_soft")

    add_panel(slide, 0.72, 2.92, 4.65, 3.88)
    add_text(slide, 0.95, 3.12, 1.8, 0.22, "Sample POD", size=12.5, bold=True, color="ink")
    add_image_contain(slide, "pod-examples/lets-transport-clean.jpg", 0.95, 3.46, 4.15, 2.78)
    add_text(slide, 0.95, 6.42, 4.0, 0.18, "Example source: mixed print + handwriting proof with visible signoff area.", size=7.8, color="ink_soft")

    add_panel(slide, 5.75, 2.92, 4.95, 1.18, fill="panel_soft")
    add_text(slide, 6.0, 3.15, 4.35, 0.28, "Extracted fields", size=15, bold=True, color="ink")
    add_text(slide, 6.0, 3.52, 4.35, 0.28, "CNS No A138687 | invoice 8530137060 | route Krishnapatnam to Thrissur", size=9.2, color="ink_soft")

    add_signal_block(slide, 5.75, 4.45, 2.35, 1.05, "Visual checks", "Signature present, stamp present, lower acknowledgement area visible")
    add_signal_block(slide, 8.35, 4.45, 2.35, 1.05, "Exception readout", "Handwritten remarks detected and routed with the extracted POD record")

    add_panel(slide, 5.75, 5.85, 4.95, 0.95, fill="navy", line="navy")
    add_text(slide, 6.0, 6.03, 4.35, 0.22, "Structured output", size=12.5, bold=True, color="white")
    add_text(slide, 6.0, 6.33, 4.4, 0.26, "{ document_type: \"POD\", signature: true, stamp: true, review_status: \"ready\" }", size=8.4, color="white")


def build_slide_4():
    slide = new_slide("Workflow Fit", 5)
    add_eyebrow(slide, 0.72, 1.0, "END-TO-END FLOW", w=1.95)
    add_text(slide, 0.72, 1.48, 8.4, 0.72, "Fracto can sit inside the {{lead_company_name}} proof and billing workflow.", size=21.5, bold=True, color="ink", font="Space Grotesk")
    steps = [
        ("01", "Capture", "{{lead_company_name}} uploads {{sample_document_types}} from drivers, hubs, or ops teams."),
        ("02", "Classify", "Each document is identified before the right logic runs."),
        ("03", "Extract and Check", "Fields are extracted and checks run for signatures, stamps, and remarks."),
        ("04", "Route", "Exceptions route to {{review_team}} while clean files move to {{downstream_system}}."),
    ]
    for idx, (no, title, body) in enumerate(steps):
        x = 0.72 + idx * 2.62
        add_panel(slide, x, 2.42, 2.35, 1.92)
        add_text(slide, x + 0.18, 2.64, 0.55, 0.35, no, size=19, bold=True, color="accent", font="Space Grotesk")
        add_text(slide, x + 0.18, 3.05, 1.9, 0.3, title, size=13.5, bold=True, color="ink")
        add_text(slide, x + 0.18, 3.43, 1.95, 0.58, body, size=9.4, color="ink_soft")
        if idx < 3:
            add_text(slide, x + 2.37, 3.12, 0.22, 0.25, ">", size=15, bold=True, color="accent", align=PP_ALIGN.CENTER)
    add_panel(slide, 0.72, 5.0, 5.1, 1.78)
    add_text(slide, 0.95, 5.22, 4.4, 0.28, "High-value insertion points", size=15.5, bold=True, color="ink")
    add_multiline(slide, 0.98, 5.62, 4.45, 0.98, [
        "{{primary_workflow}} desks", "Billing proof checks", "Hub ops and exception handling", "Trip sheet and document matching", "Customer dispute workflows"
    ], size=8.7, bullet=True)
    add_panel(slide, 6.2, 5.0, 4.55, 1.78, fill="panel_soft")
    add_text(slide, 6.43, 5.22, 3.8, 0.28, "Decision signals", size=15.5, bold=True, color="ink")
    add_multiline(slide, 6.46, 5.62, 3.8, 0.98, [
        "Signature present or missing", "Stamp present or missing", "Handwritten remarks detected", "Field confidence and exception flags"
    ], size=8.7, bullet=True)


def build_slide_5():
    slide = new_slide("Reference Set", 2)
    add_eyebrow(slide, 0.72, 1.0, "TRANSPORT OPERATIONS", w=2.25)
    add_text(slide, 0.72, 1.48, 5.75, 1.04, "Fracto supports document-heavy road transport workflows.", size=21.0, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.58, 5.65, 0.42, "For {{lead_company_name}}, the same layer can support POD checks, trip documents, proof readiness, and exception routing.", size=10.5, color="ink_soft")
    logos = [
        ("clients/navata.png", "Navata SCS"),
        ("clients/onmove-cropped.jpg", "Onmove"),
        ("clients/letstransport.png", "LetsTransport"),
        ("clients/mvikas.png", "MVikas"),
        ("clients/vishal-vtc.png", "Vishal VTC"),
        ("clients/trukker.png", "Trukker"),
    ]
    add_panel(slide, 0.72, 3.35, 5.75, 2.55, fill="panel_soft")
    add_text(slide, 0.95, 3.56, 4.1, 0.25, "Relevant transport and logistics workflows", size=13.4, bold=True, color="ink")
    for idx, (path, name) in enumerate(logos):
        x = 0.95 + (idx % 3) * 1.78
        y = 4.02 + (idx // 3) * 0.9
        add_logo_card(slide, x, y, 1.5, 0.64, path, name)
    add_panel(slide, 6.85, 1.52, 3.95, 2.22)
    add_text(slide, 7.08, 1.76, 3.42, 0.68, "Why this matters for {{lead_company_name}}", size=15.2, bold=True, color="ink")
    add_multiline(slide, 7.1, 2.58, 3.35, 0.86, [
        "High POD and proof document volume",
        "Manual checks delay billing readiness",
        "Signatures, stamps, and remarks need review",
        "Exceptions need faster routing to ops teams",
    ], size=8.9, bullet=True)
    add_panel(slide, 6.85, 4.08, 3.95, 2.02, fill="panel_soft")
    add_text(slide, 7.08, 4.32, 3.42, 0.3, "What your team can expect", size=15.8, bold=True, color="ink")
    add_multiline(slide, 7.1, 4.74, 3.35, 0.94, [
        "Cleaner proof review queues",
        "Earlier incomplete-document flags",
        "Structured outputs for downstream tools",
        "Lower manual effort across review teams",
    ], size=9.8, bullet=True)


def build_slide_6():
    slide = new_slide("Operator Impact", 6)
    add_eyebrow(slide, 0.72, 1.0, "BUSINESS OUTCOMES", w=2.1)
    add_text(slide, 0.72, 1.48, 8.55, 1.28, "The value case for {{prospect_role}} at {{lead_company_name}}.", size=19.8, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.88, 8.35, 0.32, "Help improve {{outcome_1}}, {{outcome_2}}, and {{outcome_3}} across {{primary_workflow}}.", size=10.2, color="ink_soft")
    impact = [
        ("Earlier proof readiness", "PODs can be flagged as complete or incomplete faster."),
        ("Cleaner exceptions", "Missing signoff and remarks get surfaced earlier."),
        ("One workflow", "Classification, extraction, and visual checks in one layer."),
    ]
    for idx, (title, body) in enumerate(impact):
        x = 0.72 + idx * 3.42
        add_panel(slide, x, 3.32, 3.05, 1.23)
        add_text(slide, x + 0.18, 3.52, 2.55, 0.3, title, size=14.5, bold=True, color="ink", font="Space Grotesk")
        add_text(slide, x + 0.18, 3.94, 2.55, 0.33, body, size=8.4, color="ink_soft")
    add_panel(slide, 0.72, 4.9, 4.95, 1.55)
    add_text(slide, 0.95, 5.12, 4.3, 0.28, "Common workflow bottlenecks", size=15.5, bold=True, color="ink")
    add_multiline(slide, 0.98, 5.5, 4.15, 0.72, [
        "Proof-review queues move slowly", "Missing signatures are found late", "Stamp checks are still manual", "Remarks get lost in scanned copies"
    ], size=9.7, bullet=True)
    add_panel(slide, 6.05, 4.9, 4.7, 1.55, fill="panel_soft")
    add_text(slide, 6.28, 5.12, 4.05, 0.28, "What Fracto enables", size=15.5, bold=True, color="ink")
    add_multiline(slide, 6.31, 5.5, 3.9, 0.72, [
        "Faster proof triage and routing", "Structured outputs into ops tools", "Manual review only for uncertain cases", "Better visibility into proof completeness"
    ], size=9.7, bullet=True)


def build_slide_7():
    slide = new_slide("Demo Scope", 8)
    add_eyebrow(slide, 0.72, 1.0, "RECOMMENDED NEXT STEP", w=2.55)
    add_text(slide, 0.72, 1.48, 5.7, 0.72, "Start with a {{primary_workflow}} demo for {{lead_company_name}}.", size=25, bold=True, color="ink", font="Space Grotesk")
    add_panel(slide, 0.72, 2.55, 5.2, 2.3)
    add_text(slide, 0.95, 2.8, 4.5, 0.32, "Suggested scope", size=16, bold=True, color="ink")
    add_multiline(slide, 0.98, 3.25, 4.35, 1.05, [
        "Input bundle: {{sample_document_types}}",
        "Flow: {{demo_flow}}",
        "Integration: {{integration_target}}",
        "Success: {{success_metric}}",
    ], size=10.2, bullet=True)
    add_band(slide, 0.72, 5.35, 5.2, 0.98, "Outcome", "A before-and-after view of proof review speed, completeness checks, and exception handling.")
    add_panel(slide, 6.48, 1.52, 4.25, 2.18, fill="panel_soft")
    add_text(slide, 6.72, 1.78, 3.55, 0.32, "What Fracto brings", size=16, bold=True, color="ink")
    add_multiline(slide, 6.74, 2.22, 3.35, 0.98, [
        "Document parsers for transport workflows",
        "Signature, stamp, and remark detection",
        "Review routing for edge cases",
        "Low-lift API integration",
    ], size=9.8, bullet=True)
    add_panel(slide, 6.48, 4.1, 4.25, 1.48)
    add_text(slide, 6.72, 4.34, 3.55, 0.28, "Teams this supports", size=16, bold=True, color="ink")
    cx, cy = 6.72, 4.82
    for text, width in [("Operations Head", 1.35), ("Billing Ops", 1.05), ("Proof Review", 1.2), ("Finance", 0.85), ("Engineering", 1.05)]:
        add_chip(slide, cx, cy, text, width)
        cx += width + 0.1
        if cx > 9.6:
            cx, cy = 6.72, 5.22
    add_text(slide, 6.5, 6.08, 4.15, 0.34, "A focused demo plan for {{lead_company_name}} around {{primary_workflow}} automation.", size=9.4, color="ink_soft")


def build_slide_8():
    slide = new_slide("Book Demo", 7)
    add_eyebrow(slide, 0.72, 1.0, "NEXT STEP", w=1.3)
    add_text(slide, 0.72, 1.48, 5.4, 1.55, "{{prospect_name}}, want to review this on a few {{lead_company_name}} sample files?", size=20.5, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 3.15, 5.15, 0.48, "Book a demo to review {{primary_workflow}}, document parsing, signature and stamp checks, plus exception routing on real transport files.", size=10.6, color="ink_soft")
    add_panel(slide, 0.72, 4.02, 5.2, 1.68)
    add_text(slide, 0.95, 4.24, 3.8, 0.28, "Suggested demo scope", size=14, bold=True, color="ink")
    add_multiline(slide, 0.98, 4.65, 4.45, 0.76, [
        "Input bundle: {{sample_document_types}}",
        "Flow: {{demo_flow}}",
        "Success metric: {{success_metric}}",
    ], size=9.2, bullet=True)
    add_band(slide, 0.72, 5.9, 5.2, 1.05, "Demo Focus", "Review real files first. Then scope the right workflow.")
    add_panel(slide, 6.45, 1.28, 4.3, 5.55, fill="navy", line="navy")
    add_eyebrow(slide, 6.75, 1.68, "TALK TO FRACTO", w=1.65)
    add_text(slide, 6.75, 2.25, 3.55, 0.65, "Book a demo and review capabilities.", size=22, bold=True, color="white", font="Space Grotesk")
    contact = [("Email", "sales@fracto.tech"), ("Phone", "+91 84483 33277"), ("Website", "fracto.tech")]
    y = 3.22
    for label, value in contact:
        add_panel(slide, 6.78, y, 3.65, 0.58, fill="accent_dark", line="accent_dark")
        add_text(slide, 6.95, y + 0.08, 0.9, 0.18, label.upper(), size=7.5, bold=True, color="white")
        add_text(slide, 7.82, y + 0.06, 2.35, 0.22, value, size=12, bold=True, color="white", font="Space Grotesk")
        y += 0.78
    add_text(slide, 6.78, 5.92, 3.6, 0.26, "Next step: review a small {{primary_workflow}} sample set.", size=8.8, color="white")


for builder in [
    build_slide_1,
    build_slide_5,
    build_slide_2,
    build_slide_4_visual_proof,
    build_slide_4,
    build_slide_6,
    build_slide_8,
]:
    builder()

props = prs.core_properties
props.title = "Fracto Road Transport Editable Dynamic Sales Deck"
props.subject = "Personalized client-facing sales deck with editable variables"
props.author = "Fracto"
TMP_OUT = OUT.with_suffix(".tmp.pptx")
if TMP_OUT.exists():
    TMP_OUT.unlink()
prs.save(TMP_OUT)
TMP_OUT.replace(OUT)
print(OUT)
