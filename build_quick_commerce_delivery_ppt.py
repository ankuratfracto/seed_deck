from pathlib import Path
import re
from PIL import Image, ImageDraw, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "quick-commerce-delivery-dynamic.pptx"


COLORS = {
    "ink": RGBColor(20, 33, 61),
    "ink_soft": RGBColor(50, 68, 99),
    "paper": RGBColor(255, 253, 248),
    "panel": RGBColor(255, 255, 255),
    "panel_soft": RGBColor(245, 247, 251),
    "line": RGBColor(215, 222, 234),
    "accent": RGBColor(14, 143, 125),
    "accent_dark": RGBColor(11, 109, 96),
    "accent_soft": RGBColor(218, 246, 240),
    "signal": RGBColor(242, 163, 59),
    "signal_soft": RGBColor(255, 241, 219),
    "berry": RGBColor(191, 91, 118),
    "navy": RGBColor(21, 49, 77),
    "white": RGBColor(255, 255, 255),
}


prs = Presentation()
prs.slide_width = Inches(11.69)
prs.slide_height = Inches(8.27)
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["paper"]


def set_run(run, size=12, bold=False, color="ink", font="IBM Plex Sans"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]


def add_text(slide, x, y, w, h, text, size=12, bold=False, color="ink", align=PP_ALIGN.LEFT, font="IBM Plex Sans"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, color=color, font=font)
    return shape


def add_multiline(slide, x, y, w, h, lines, size=12, color="ink_soft", bullet=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {line}" if bullet else line
        p.space_after = Pt(5)
        for run in p.runs:
            set_run(run, size=size, color=color)
    return shape


def add_panel(slide, x, y, w, h, fill="panel", line="line", radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    shape.line.width = Pt(1)
    shape.adjustments[0] = radius
    return shape


def add_band(slide, x, y, w, h, title, body):
    shape = add_panel(slide, x, y, w, h, fill="navy", line="navy")
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.35, title, size=18, bold=True, color="white", font="Space Grotesk")
    add_text(slide, x + 0.18, y + 0.58, w - 0.36, h - 0.7, body, size=11.2, color="white")
    return shape


def add_tag(slide, x, y, text, w=1.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["signal_soft"]
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    add_text(slide, x + 0.05, y + 0.055, w - 0.1, 0.18, text, size=7.8, bold=True, color="accent_dark", align=PP_ALIGN.CENTER)
    return shape


def add_eyebrow(slide, x, y, text, w=2.6):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["accent_soft"]
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    add_text(slide, x + 0.12, y + 0.07, w - 0.24, 0.2, text, size=8.4, bold=True, color="accent_dark", align=PP_ALIGN.CENTER)
    return shape


def add_chip(slide, x, y, text, w):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["panel"]
    shape.line.color.rgb = COLORS["line"]
    shape.adjustments[0] = 0.5
    add_text(slide, x + 0.05, y + 0.075, w - 0.1, 0.18, text, size=8.5, bold=True, color="ink_soft", align=PP_ALIGN.CENTER)
    return shape


def add_card(slide, x, y, w, h, tag, title, body, tag_w=0.72):
    add_panel(slide, x, y, w, h)
    add_tag(slide, x + 0.18, y + 0.16, tag, w=tag_w)
    add_text(slide, x + 0.18, y + 0.58, w - 0.36, 0.42, title, size=12.2, bold=True, color="ink")
    add_text(slide, x + 0.18, y + 1.06, w - 0.36, h - 1.14, body, size=8.6, color="ink_soft")


def add_metric(slide, x, y, w, h, value, label):
    add_panel(slide, x, y, w, h, fill="panel")
    add_text(slide, x + 0.15, y + 0.16, w - 0.3, 0.54, value, size=11.6, bold=True, color="accent_dark")
    add_text(slide, x + 0.15, y + 0.76, w - 0.3, h - 0.88, label, size=7.9, color="ink_soft")


def add_header(slide, section):
    if (ROOT / "fracto_logo.png").exists():
        slide.shapes.add_picture(str(ROOT / "fracto_logo.png"), Inches(0.55), Inches(0.28), width=Inches(0.72), height=Inches(0.225))
    add_text(slide, 9.1, 0.36, 2.0, 0.25, section, size=8.3, bold=True, color="ink_soft", align=PP_ALIGN.RIGHT)


def add_footer(slide, page):
    add_text(slide, 0.55, 7.78, 3.2, 0.24, "Prepared for {{lead_company_name}}", size=8.5, color="ink_soft")
    add_text(slide, 10.55, 7.78, 0.55, 0.24, f"{page:02d}", size=8.5, color="ink_soft", align=PP_ALIGN.RIGHT)


def new_slide(section, page):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_header(slide, section)
    add_footer(slide, page)
    return slide


def add_logo_card(slide, x, y, w, h, img_path, fallback):
    dark_logo = "qwqer" in img_path.lower()
    add_panel(slide, x, y, w, h, fill="navy" if dark_logo else "panel", line="line")
    path = ROOT / img_path
    if path.exists():
        with Image.open(path) as img:
            img_w, img_h = img.size
        pad_x, pad_y = 0.12, 0.1
        box_w, box_h = w - (pad_x * 2), h - (pad_y * 2)
        img_ratio = img_w / img_h
        box_ratio = box_w / box_h
        if img_ratio > box_ratio:
            out_w = box_w
            out_h = box_w / img_ratio
        else:
            out_h = box_h
            out_w = box_h * img_ratio
        out_x = x + (w - out_w) / 2
        out_y = y + (h - out_h) / 2
        slide.shapes.add_picture(str(path), Inches(out_x), Inches(out_y), width=Inches(out_w), height=Inches(out_h))
    else:
        add_text(slide, x + 0.15, y + 0.33, w - 0.3, 0.26, fallback, size=10.5, bold=True, color="white" if dark_logo else "ink_soft", align=PP_ALIGN.CENTER)


def add_image_contain(slide, img_path, x, y, w, h):
    path = ROOT / img_path
    with Image.open(path) as img:
        img_w, img_h = img.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        out_w = w
        out_h = w / img_ratio
    else:
        out_h = h
        out_w = h * img_ratio
    out_x = x + (w - out_w) / 2
    out_y = y + (h - out_h) / 2
    return slide.shapes.add_picture(str(path), Inches(out_x), Inches(out_y), width=Inches(out_w), height=Inches(out_h))


def add_signal_block(slide, x, y, w, h, label, value):
    add_panel(slide, x, y, w, h, fill="panel")
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.22, label, size=8.2, bold=True, color="accent_dark")
    add_text(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.58, value, size=9.2, color="ink_soft")


def add_field_row(slide, x, y, label, value, w=2.9):
    add_text(slide, x, y, 0.74, 0.18, label.upper(), size=5.4, bold=True, color="ink_soft")
    add_text(slide, x + 0.78, y - 0.01, w - 0.78, 0.2, value, size=7.5, bold=True, color="ink")


def add_sample_document(slide, x, y, w, h, tag, title, subtitle, rows, fill="panel"):
    add_panel(slide, x, y, w, h, fill=fill)
    add_tag(slide, x + 0.16, y + 0.14, tag, w=0.72)
    add_text(slide, x + 0.98, y + 0.15, w - 1.14, 0.22, title, size=11.4, bold=True, color="ink")
    add_text(slide, x + 0.16, y + 0.49, w - 0.32, 0.2, subtitle, size=7.2, color="ink_soft")
    cy = y + 0.82
    for label, value in rows:
        add_field_row(slide, x + 0.2, cy, label, value, w=w - 0.4)
        cy += 0.3


def add_compact_json(slide, x, y, w, h, title, lines):
    add_panel(slide, x, y, w, h, fill="navy", line="navy")
    add_text(slide, x + 0.18, y + 0.15, w - 0.36, 0.2, title, size=9.8, bold=True, color="white")
    add_multiline(slide, x + 0.18, y + 0.45, w - 0.36, h - 0.54, lines, size=7.0, color="white", bullet=False)


def ensure_parcel_photo():
    out = ROOT / "quick-commerce-parcel-delivery-sample.png"
    img = Image.new("RGB", (1100, 760), (236, 238, 235))
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 515, 1100, 760), fill=(214, 217, 211))
    draw.rectangle((0, 0, 1100, 515), fill=(246, 246, 242))
    draw.rounded_rectangle((120, 140, 880, 620), radius=24, fill=(231, 186, 126), outline=(170, 122, 71), width=5)
    draw.line((120, 250, 880, 250), fill=(158, 109, 59), width=5)
    draw.line((500, 140, 500, 620), fill=(158, 109, 59), width=4)
    draw.rounded_rectangle((575, 308, 790, 432), radius=10, fill=(250, 250, 245), outline=(57, 70, 95), width=3)
    draw.text((600, 332), "ORDER", fill=(21, 49, 77))
    draw.text((600, 368), "QC-49281", fill=(21, 49, 77))
    draw.rectangle((605, 402, 760, 412), fill=(21, 49, 77))
    draw.rectangle((620, 417, 720, 424), fill=(21, 49, 77))
    draw.rounded_rectangle((55, 54, 480, 112), radius=16, fill=(21, 49, 77))
    draw.text((84, 73), "RIDER APP UPLOAD", fill=(255, 255, 255))
    draw.rounded_rectangle((735, 548, 1045, 610), radius=14, fill=(218, 246, 240), outline=(14, 143, 125), width=3)
    draw.text((760, 568), "Package visible | Label visible", fill=(11, 109, 96))
    img.save(out)
    return out


def ensure_qwqer_parcel_photo():
    src = Path("/Users/ankurgupta/Desktop/a77d001c20204cfc_Image_20260511_222147_10301726666613946.jpg")
    out = ROOT / "quick-commerce-qwqer-parcel-visible.jpg"
    if not src.exists():
        return ensure_parcel_photo()
    img = ImageOps.exif_transpose(Image.open(src)).convert("RGB")
    img.save(out, quality=92, optimize=True)
    return out


def build_slide_1():
    slide = new_slide("Quick Commerce Ops Deck", 1)
    add_eyebrow(slide, 0.72, 1.02, "AI INFRASTRUCTURE FOR RIDER AND PARCEL OPS", w=4.25)
    add_text(slide, 0.72, 1.58, 6.15, 1.28, "Automate {{primary_workflow}} for {{lead_company_name}}.", size=25, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.92, 6.05, 0.72, "Prepared for {{prospect_name}} and the {{lead_company_name}} team. Fracto can verify rider onboarding documents, classify delivery images, and return parcel-proof decisions directly into ops workflows.", size=11.4, color="ink_soft")
    add_panel(slide, 0.72, 3.88, 5.45, 1.22)
    add_multiline(slide, 0.93, 4.07, 4.95, 0.82, [
        "Rider onboarding for Aadhaar, PAN, DL, and ID files",
        "Parcel delivery image verification",
        "POD, AWB, invoice, and manifest extraction",
        "APIs for rider apps, seller portals, and ops queues",
    ], size=9.6, bullet=True)
    for x, y, text, width in [
        (0.72, 5.35, "{{primary_workflow}}", 1.8),
        (2.67, 5.35, "Rider KYC", 1.05),
        (3.88, 5.35, "Parcel proof", 1.2),
        (0.72, 5.78, "POD / AWB", 1.05),
        (1.93, 5.78, "For {{prospect_role}}", 1.7),
    ]:
        add_chip(slide, x, y, text, width)
    add_panel(slide, 7.0, 1.02, 3.95, 4.72, fill="panel_soft")
    add_text(slide, 7.27, 1.26, 3.25, 0.58, "What this changes for {{lead_company_name}}", size=14.5, bold=True, color="ink")
    add_metric(slide, 7.27, 2.02, 1.55, 1.32, "Rider ready", "KYC fields and onboarding documents become structured")
    add_metric(slide, 9.05, 2.02, 1.55, 1.32, "Parcel proof", "Images return parcel, not parcel, or review decisions")
    add_metric(slide, 7.27, 3.5, 1.55, 1.32, "Docs ready", "POD, AWB, invoice, and manifest extraction")
    add_metric(slide, 9.05, 3.5, 1.55, 1.32, "Faster ops", "Exceptions route before customer or billing delays")
    add_band(slide, 7.05, 5.8, 3.9, 1.25, "Positioning", "One verification layer for {{lead_company_name}} across rider documents and parcel proof.")


def build_slide_2_references():
    slide = new_slide("Reference Set", 2)
    add_eyebrow(slide, 0.72, 1.0, "DELIVERY OPERATIONS", w=2.25)
    add_text(slide, 0.72, 1.48, 5.8, 1.04, "Fracto supports parcel, delivery, and e-commerce operations workflows.", size=21.0, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.58, 5.65, 0.42, "For {{lead_company_name}}, the same layer can support rider onboarding, delivery image verification, shipment documents, and exception routing.", size=10.5, color="ink_soft")
    logos = [
        ("clients/qwqer-dark.png", "Qwqer"),
        ("clients/parcelx.png", "ParcelX"),
        ("clients/jeebly.png", "Jeebly"),
        ("clients/letstransport.png", "LetsTransport"),
        ("clients/fleetx.png", "Fleetx"),
        ("clients/petpooja.png", "Petpooja"),
    ]
    add_panel(slide, 0.72, 3.35, 5.75, 2.55, fill="panel_soft")
    add_text(slide, 0.95, 3.56, 4.1, 0.25, "Relevant delivery and e-commerce workflows", size=13.4, bold=True, color="ink")
    for idx, (path, name) in enumerate(logos):
        x = 0.95 + (idx % 3) * 1.78
        y = 4.02 + (idx // 3) * 0.9
        add_logo_card(slide, x, y, 1.5, 0.64, path, name)
    add_panel(slide, 6.85, 1.52, 3.95, 2.22)
    add_text(slide, 7.08, 1.76, 3.42, 0.68, "Why this matters for {{lead_company_name}}", size=15.2, bold=True, color="ink")
    add_multiline(slide, 7.1, 2.58, 3.35, 0.86, [
        "High daily proof-image volume",
        "Rider and seller uploads vary in quality",
        "Parcel proof decisions need to be fast",
        "Exceptions need quick routing to ops teams",
    ], size=8.9, bullet=True)
    add_panel(slide, 6.85, 4.08, 3.95, 2.02, fill="panel_soft")
    add_text(slide, 7.08, 4.32, 3.42, 0.3, "What your team can expect", size=15.8, bold=True, color="ink")
    add_multiline(slide, 7.1, 4.74, 3.35, 0.94, [
        "Automated parcel / not-parcel flags",
        "Cleaner rider onboarding records",
        "Structured shipment document outputs",
        "Lower manual review across ops queues",
    ], size=9.8, bullet=True)


def build_slide_3_problem():
    slide = new_slide("Why Delivery Ops Break Down", 3)
    add_eyebrow(slide, 0.72, 1.0, "THE OPERATIONAL PROBLEM", w=2.35)
    add_text(slide, 0.72, 1.48, 6.15, 0.98, "Delivery teams move fast, but image proof and onboarding checks still get stuck in queues.", size=20.5, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.58, 5.95, 0.48, "For {{lead_company_name}}, {{pain_point_1}}, {{pain_point_2}}, and {{pain_point_3}} can slow down {{primary_workflow}}.", size=10.6, color="ink_soft")
    cards = [
        ("Images", "Rider uploads are inconsistent", "Ops teams still check whether a photo actually shows a parcel."),
        ("KYC", "Rider onboarding needs cleanup", "Aadhaar, PAN, DL, and ID uploads arrive blurred or incomplete."),
        ("Docs", "Shipment files are copied manually", "POD, AWB, invoice, and manifest fields still move across tools by hand."),
        ("Scale", "Exceptions pile up quickly", "Daily delivery spikes create review backlogs across hubs and ops teams."),
    ]
    coords = [(0.72, 3.28), (3.55, 3.28), (0.72, 5.1), (3.55, 5.1)]
    for (tag, title, body), (x, y) in zip(cards, coords):
        add_card(slide, x, y, 2.55, 1.58, tag, title, body, tag_w=1.0)
    add_panel(slide, 6.85, 1.52, 3.95, 2.2)
    add_text(slide, 7.08, 1.78, 3.42, 0.33, "Typical failure pattern", size=17, bold=True, color="ink")
    add_multiline(slide, 7.1, 2.24, 3.35, 1.05, [
        "Random, duplicate, blurry, or cropped delivery photos",
        "Manual checks for package visibility and label presence",
        "Rider KYC and shipment files reviewed in separate tools",
        "Customer, billing, and claims teams wait on proof status",
    ], size=9.6, bullet=True)
    add_panel(slide, 6.85, 4.1, 3.95, 2.12, fill="panel_soft")
    add_text(slide, 7.08, 4.35, 3.42, 0.33, "What teams need instead", size=17, bold=True, color="ink")
    add_multiline(slide, 7.1, 4.8, 3.35, 1.0, [
        "Automated valid or invalid parcel image decisions",
        "Extraction for onboarding and shipment documents",
        "Review only uncertain cases",
        "API outputs for rider apps and ops tools",
    ], size=9.8, bullet=True)


def build_slide_4_rider_onboarding():
    slide = new_slide("Rider Onboarding", 4)
    add_eyebrow(slide, 0.72, 1.0, "RIDER KYC + VEHICLE CHECKS", w=2.95)
    add_text(slide, 0.72, 1.46, 8.9, 0.62, "Rider onboarding: verify identity, documents, and vehicle before first delivery.", size=19.4, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.15, 8.65, 0.32, "The same Fracto layer can parse rider files, mask Aadhaar, validate DL/RC data, and return a clean onboarding decision.", size=9.8, color="ink_soft")

    add_panel(slide, 0.72, 2.72, 4.95, 4.22)
    add_text(slide, 0.95, 2.95, 3.55, 0.25, "Sample onboarding uploads", size=13.8, bold=True, color="ink")
    add_sample_document(slide, 0.95, 3.42, 2.15, 2.15, "KYC", "Aadhaar", "Masked identity image", [
        ("Name", "RAVI KUMAR"),
        ("Aadhaar", "XXXX XXXX 4321"),
        ("DOB", "19XX"),
        ("Address", "Bengaluru"),
    ], fill="panel_soft")
    add_sample_document(slide, 3.25, 3.42, 2.15, 2.15, "PAN", "PAN Card", "Tax ID image", [
        ("Name", "RAVI KUMAR"),
        ("PAN", "ABCDE1234F"),
        ("Status", "Name match"),
        ("DOB", "19XX"),
    ])
    add_panel(slide, 0.95, 5.88, 4.45, 0.66, fill="panel_soft")
    add_text(slide, 1.14, 6.05, 1.42, 0.18, "Verification inputs", size=8.6, bold=True, color="accent_dark")
    add_text(slide, 2.48, 6.05, 2.7, 0.18, "DL: KA05 20XX 1234567 | RC: KA-05-AB-1234", size=7.4, color="ink_soft")

    add_panel(slide, 5.95, 2.72, 2.25, 2.08, fill="panel_soft")
    add_text(slide, 6.16, 2.96, 1.78, 0.25, "Extracted data", size=13.2, bold=True, color="ink")
    add_multiline(slide, 6.18, 3.36, 1.75, 1.08, [
        "rider_name: Ravi Kumar",
        "aadhaar: XXXX-XXXX-4321",
        "pan: ABCDE1234F",
        "dl_no: KA05 20XX 1234567",
        "rc_no: KA-05-AB-1234",
    ], size=7.4, bullet=False)

    add_panel(slide, 8.42, 2.72, 2.25, 2.08)
    add_text(slide, 8.63, 2.96, 1.72, 0.25, "Fracto signals", size=13.2, bold=True, color="ink")
    add_multiline(slide, 8.65, 3.36, 1.72, 1.08, [
        "aadhaar_masked: true",
        "pan_name_match: true",
        "dl_status: active",
        "rc_status: verified",
        "route_to: {{downstream_system}}",
    ], size=7.4, bullet=False)

    add_compact_json(slide, 5.95, 5.08, 2.25, 1.55, "Parivahan DL API", [
        "status: active",
        "holder_match: true",
        "vehicle_class: 2W/LMV",
        "valid_till: 2028-12",
    ])
    add_compact_json(slide, 8.42, 5.08, 2.25, 1.55, "Parivahan RC API", [
        "status: verified",
        "owner_match: true",
        "fitness: active",
        "insurance: valid",
    ])


def build_slide_5_visual_proof():
    slide = new_slide("Visual Proof", 5)
    img_path = ensure_qwqer_parcel_photo()
    add_eyebrow(slide, 0.72, 1.0, "PARCEL IMAGE + DECISION", w=2.55)
    add_text(slide, 0.72, 1.46, 8.95, 0.62, "Qwqer use case: confirm whether a parcel is visible in the delivery image.", size=19.4, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.15, 8.75, 0.32, "Fracto reads the rider-uploaded image and returns a simple decision that can move directly into the delivery workflow.", size=9.8, color="ink_soft")
    add_panel(slide, 0.72, 2.72, 4.05, 4.35)
    add_text(slide, 0.95, 2.95, 3.35, 0.25, "Rider delivery image", size=13.8, bold=True, color="ink")
    add_image_contain(slide, img_path.name, 0.95, 3.28, 3.58, 3.25)
    add_text(slide, 0.95, 6.66, 3.4, 0.18, "Parcel and shipping label are visible in the uploaded proof image.", size=7.8, color="ink_soft")

    add_panel(slide, 5.05, 2.72, 5.7, 2.08, fill="panel_soft")
    add_text(slide, 5.32, 2.96, 4.6, 0.28, "Extracted fields", size=15.2, bold=True, color="ink")
    add_signal_block(slide, 5.32, 3.4, 1.85, 0.82, "Is Parcel", "true")
    add_signal_block(slide, 7.38, 3.4, 3.0, 0.82, "Description", "Brown package with shipping label and branded tape.")

    add_panel(slide, 5.05, 5.05, 2.7, 1.62)
    add_text(slide, 5.28, 5.27, 2.2, 0.28, "Delivery decision", size=14.2, bold=True, color="ink")
    add_multiline(slide, 5.3, 5.66, 2.05, 0.7, [
        "parcel_visible: true",
        "proof_status: accept",
        "review_required: false",
    ], size=8.2, bullet=False)
    add_compact_json(slide, 8.05, 5.05, 2.7, 1.62, "Ops Callback", [
        "order_id: {{order_id}}",
        "image_type: parcel",
        "route_to: {{downstream_system}}",
        "send_review: false",
    ])


def build_slide_6_workflow():
    slide = new_slide("Workflow Fit", 6)
    add_eyebrow(slide, 0.72, 1.0, "END-TO-END FLOW", w=1.95)
    add_text(slide, 0.72, 1.48, 8.4, 0.72, "Fracto can sit inside the {{lead_company_name}} rider and parcel workflow.", size=21.5, bold=True, color="ink", font="Space Grotesk")
    steps = [
        ("01", "Capture", "Riders, sellers, or hubs upload KYC files, parcel photos, and shipment documents."),
        ("02", "Classify", "Fracto identifies whether each upload is KYC, parcel image, POD, AWB, invoice, or manifest."),
        ("03", "Verify", "Images are checked for parcel/not-parcel status while documents are parsed into fields."),
        ("04", "Route", "Clean decisions move to {{downstream_system}} while exceptions route to {{review_team}}."),
    ]
    for idx, (no, title, body) in enumerate(steps):
        x = 0.72 + idx * 2.62
        add_panel(slide, x, 2.42, 2.35, 1.92)
        add_text(slide, x + 0.18, 2.64, 0.55, 0.35, no, size=19, bold=True, color="accent", font="Space Grotesk")
        add_text(slide, x + 0.18, 3.05, 1.9, 0.3, title, size=13.5, bold=True, color="ink")
        add_text(slide, x + 0.18, 3.43, 1.95, 0.58, body, size=8.6, color="ink_soft")
        if idx < 3:
            add_text(slide, x + 2.37, 3.12, 0.22, 0.25, ">", size=15, bold=True, color="accent", align=PP_ALIGN.CENTER)
    add_panel(slide, 0.72, 5.0, 5.1, 1.78)
    add_text(slide, 0.95, 5.22, 4.4, 0.28, "High-value insertion points", size=15.5, bold=True, color="ink")
    add_multiline(slide, 0.98, 5.62, 4.45, 0.98, [
        "Rider onboarding",
        "Delivery completion photos",
        "Return and claim evidence",
        "Seller or hub document uploads",
        "Billing proof workflows",
    ], size=8.7, bullet=True)
    add_panel(slide, 6.2, 5.0, 4.55, 1.78, fill="panel_soft")
    add_text(slide, 6.43, 5.22, 3.8, 0.28, "Decision signals", size=15.5, bold=True, color="ink")
    add_multiline(slide, 6.46, 5.62, 3.8, 0.98, [
        "Parcel, not parcel, or review",
        "Package and label visibility",
        "KYC and shipment field confidence",
        "Duplicate, blurred, or irrelevant uploads",
    ], size=8.7, bullet=True)


def build_slide_7_impact():
    slide = new_slide("Operator Impact", 7)
    add_eyebrow(slide, 0.72, 1.0, "BUSINESS OUTCOMES", w=2.1)
    add_text(slide, 0.72, 1.48, 8.55, 1.28, "The value case for {{prospect_role}} at {{lead_company_name}}.", size=19.8, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.88, 8.35, 0.32, "Help improve {{outcome_1}}, {{outcome_2}}, and {{outcome_3}} across {{primary_workflow}}.", size=10.2, color="ink_soft")
    impact = [
        ("Faster rider activation", "KYC and onboarding documents move to decision faster."),
        ("Cleaner delivery proof", "Parcel photos can be accepted, rejected, or routed to review."),
        ("One proof layer", "Images, rider files, PODs, AWBs, invoices, and manifests in one API."),
    ]
    for idx, (title, body) in enumerate(impact):
        x = 0.72 + idx * 3.42
        add_panel(slide, x, 3.32, 3.05, 1.23)
        add_text(slide, x + 0.18, 3.52, 2.55, 0.3, title, size=14.5, bold=True, color="ink", font="Space Grotesk")
        add_text(slide, x + 0.18, 3.94, 2.55, 0.33, body, size=8.2, color="ink_soft")
    add_panel(slide, 0.72, 4.9, 4.95, 1.55)
    add_text(slide, 0.95, 5.12, 4.3, 0.28, "Common workflow bottlenecks", size=15.5, bold=True, color="ink")
    add_multiline(slide, 0.98, 5.5, 4.15, 0.72, [
        "Rider photos are checked manually",
        "Random or unusable images enter ops queues",
        "KYC and shipment fields are retyped",
        "Claims and billing wait for proof status",
    ], size=9.4, bullet=True)
    add_panel(slide, 6.05, 4.9, 4.7, 1.55, fill="panel_soft")
    add_text(slide, 6.28, 5.12, 4.05, 0.28, "What Fracto enables", size=15.5, bold=True, color="ink")
    add_multiline(slide, 6.31, 5.5, 3.9, 0.72, [
        "Automated image decisioning",
        "Structured rider and shipment outputs",
        "Reviewer queues only for exceptions",
        "Better delivery proof and audit trail",
    ], size=9.4, bullet=True)


def build_slide_8_demo():
    slide = new_slide("Book Demo", 8)
    add_eyebrow(slide, 0.72, 1.0, "NEXT STEP", w=1.3)
    add_text(slide, 0.72, 1.48, 5.4, 1.55, "{{prospect_name}}, want to review this on a few {{lead_company_name}} delivery files?", size=20.5, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 3.15, 5.15, 0.48, "Book a demo to review rider onboarding documents, parcel image verification, shipment document parsing, and exception routing on sample files.", size=10.6, color="ink_soft")
    add_panel(slide, 0.72, 4.02, 5.2, 1.68)
    add_text(slide, 0.95, 4.24, 3.8, 0.28, "Suggested demo scope", size=14, bold=True, color="ink")
    add_multiline(slide, 0.98, 4.65, 4.45, 0.76, [
        "Input bundle: {{sample_document_types}}",
        "Flow: {{demo_flow}}",
        "Success metric: {{success_metric}}",
    ], size=9.2, bullet=True)
    add_band(slide, 0.72, 5.9, 5.2, 1.05, "Demo Focus", "Review real delivery images first. Then scope the right integration path.")
    add_panel(slide, 6.45, 1.28, 4.3, 5.55, fill="navy", line="navy")
    add_eyebrow(slide, 6.75, 1.68, "TALK TO FRACTO", w=1.65)
    add_text(slide, 6.75, 2.25, 3.55, 0.65, "Book a demo and review capabilities.", size=22, bold=True, color="white", font="Space Grotesk")
    contact = [("Email", "sales@fracto.tech"), ("Phone", "+91 84483 33277"), ("Website", "fracto.tech")]
    y = 3.22
    for label, value in contact:
        add_panel(slide, 6.78, y, 3.65, 0.58, fill="accent_dark", line="accent_dark")
        add_text(slide, 6.95, y + 0.08, 0.9, 0.18, label.upper(), size=7.5, bold=True, color="white")
        add_text(slide, 7.82, y + 0.06, 2.35, 0.22, value, size=12, bold=True, color="white", font="Space Grotesk")
        y += 0.78
    add_text(slide, 6.78, 5.92, 3.6, 0.26, "Next step: review a small {{primary_workflow}} sample set.", size=8.8, color="white")


for builder in [
    build_slide_1,
    build_slide_2_references,
    build_slide_3_problem,
    build_slide_4_rider_onboarding,
    build_slide_5_visual_proof,
    build_slide_6_workflow,
    build_slide_7_impact,
    build_slide_8_demo,
]:
    builder()

props = prs.core_properties
props.title = "Fracto Quick Commerce Delivery Dynamic Sales Deck"
props.subject = "Personalized client-facing sales deck with editable variables"
props.author = "Fracto"
TMP_OUT = OUT.with_suffix(".tmp.pptx")
if TMP_OUT.exists():
    TMP_OUT.unlink()
prs.save(TMP_OUT)
TMP_OUT.replace(OUT)
print(OUT)
