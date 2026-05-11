from pathlib import Path
import re
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "mobility-driver-onboarding-dynamic.pptx"


COLORS = {
    "ink": RGBColor(20, 33, 61),
    "ink_soft": RGBColor(50, 68, 99),
    "paper": RGBColor(255, 253, 248),
    "panel": RGBColor(255, 255, 255),
    "panel_soft": RGBColor(245, 247, 251),
    "line": RGBColor(215, 222, 234),
    "accent": RGBColor(14, 143, 125),
    "accent_dark": RGBColor(11, 109, 96),
    "accent_soft": RGBColor(218, 246, 240),
    "signal": RGBColor(242, 163, 59),
    "signal_soft": RGBColor(255, 241, 219),
    "berry": RGBColor(191, 91, 118),
    "navy": RGBColor(21, 49, 77),
    "white": RGBColor(255, 255, 255),
}


prs = Presentation()
prs.slide_width = Inches(11.69)
prs.slide_height = Inches(8.27)
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["paper"]


def set_run(run, size=12, bold=False, color="ink", font="IBM Plex Sans"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]


def add_text(slide, x, y, w, h, text, size=12, bold=False, color="ink", align=PP_ALIGN.LEFT, font="IBM Plex Sans"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, color=color, font=font)
    return shape


def add_multiline(slide, x, y, w, h, lines, size=12, color="ink_soft", bullet=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {line}" if bullet else line
        p.space_after = Pt(5)
        for run in p.runs:
            set_run(run, size=size, color=color)
    return shape


def add_panel(slide, x, y, w, h, fill="panel", line="line", radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    shape.line.width = Pt(1)
    shape.adjustments[0] = radius
    return shape


def add_band(slide, x, y, w, h, title, body):
    shape = add_panel(slide, x, y, w, h, fill="navy", line="navy")
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.35, title, size=18, bold=True, color="white", font="Space Grotesk")
    add_text(slide, x + 0.18, y + 0.58, w - 0.36, h - 0.7, body, size=11.2, color="white")
    return shape


def add_tag(slide, x, y, text, w=1.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["signal_soft"]
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    add_text(slide, x + 0.05, y + 0.055, w - 0.1, 0.18, text, size=7.8, bold=True, color="accent_dark", align=PP_ALIGN.CENTER)
    return shape


def add_eyebrow(slide, x, y, text, w=2.6):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["accent_soft"]
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    add_text(slide, x + 0.12, y + 0.07, w - 0.24, 0.2, text, size=8.4, bold=True, color="accent_dark", align=PP_ALIGN.CENTER)
    return shape


def add_chip(slide, x, y, text, w):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["panel"]
    shape.line.color.rgb = COLORS["line"]
    shape.adjustments[0] = 0.5
    add_text(slide, x + 0.05, y + 0.075, w - 0.1, 0.18, text, size=8.5, bold=True, color="ink_soft", align=PP_ALIGN.CENTER)
    return shape


def add_card(slide, x, y, w, h, tag, title, body, tag_w=0.72):
    add_panel(slide, x, y, w, h)
    add_tag(slide, x + 0.18, y + 0.16, tag, w=tag_w)
    add_text(slide, x + 0.18, y + 0.58, w - 0.36, 0.42, title, size=12.2, bold=True, color="ink")
    add_text(slide, x + 0.18, y + 1.06, w - 0.36, h - 1.14, body, size=8.6, color="ink_soft")


def add_metric(slide, x, y, w, h, value, label):
    add_panel(slide, x, y, w, h, fill="panel")
    add_text(slide, x + 0.15, y + 0.16, w - 0.3, 0.54, value, size=11.6, bold=True, color="accent_dark")
    add_text(slide, x + 0.15, y + 0.76, w - 0.3, h - 0.88, label, size=7.9, color="ink_soft")


def add_header(slide, section):
    if (ROOT / "fracto_logo.png").exists():
        slide.shapes.add_picture(str(ROOT / "fracto_logo.png"), Inches(0.55), Inches(0.28), width=Inches(0.72), height=Inches(0.225))
    add_text(slide, 9.1, 0.36, 2.0, 0.25, section, size=8.3, bold=True, color="ink_soft", align=PP_ALIGN.RIGHT)


def add_footer(slide, page):
    add_text(slide, 0.55, 7.78, 3.2, 0.24, "Prepared for {{lead_company_name}}", size=8.5, color="ink_soft")
    add_text(slide, 10.55, 7.78, 0.55, 0.24, f"{page:02d}", size=8.5, color="ink_soft", align=PP_ALIGN.RIGHT)


def new_slide(section, page):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_header(slide, section)
    add_footer(slide, page)
    return slide


def add_logo_card(slide, x, y, w, h, img_path, fallback):
    add_panel(slide, x, y, w, h)
    path = ROOT / img_path
    if path.exists():
        with Image.open(path) as img:
            img_w, img_h = img.size
        pad_x, pad_y = 0.12, 0.1
        box_w, box_h = w - (pad_x * 2), h - (pad_y * 2)
        img_ratio = img_w / img_h
        box_ratio = box_w / box_h
        if img_ratio > box_ratio:
            out_w = box_w
            out_h = box_w / img_ratio
        else:
            out_h = box_h
            out_w = box_h * img_ratio
        out_x = x + (w - out_w) / 2
        out_y = y + (h - out_h) / 2
        slide.shapes.add_picture(str(path), Inches(out_x), Inches(out_y), width=Inches(out_w), height=Inches(out_h))
    else:
        add_text(slide, x + 0.15, y + 0.33, w - 0.3, 0.26, fallback, size=10.5, bold=True, color="ink_soft", align=PP_ALIGN.CENTER)


def build_slide_1():
    slide = new_slide("Mobility Onboarding Deck", 1)
    add_eyebrow(slide, 0.72, 1.02, "AI INFRASTRUCTURE FOR MOBILITY ONBOARDING", w=4.0)
    add_text(slide, 0.72, 1.58, 5.95, 1.28, "Automate {{primary_workflow}} for {{lead_company_name}}.", size=25, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.92, 5.95, 0.72, "Prepared for {{prospect_name}} and the {{lead_company_name}} team. Fracto can parse Aadhaar, PAN, DL, and RC uploads, mask sensitive identity data, and validate vehicle records through Parivahan-powered checks.", size=11.4, color="ink_soft")
    add_panel(slide, 0.72, 3.88, 5.45, 1.22)
    add_multiline(slide, 0.93, 4.07, 4.95, 0.82, [
        "Aadhaar parsing with masking",
        "PAN, Driving License, and RC extraction",
        "Parivahan validation responses",
        "APIs for apps, ops desks, and reviewer queues"
    ], size=9.8, bullet=True)
    for x, y, text, width in [
        (0.72, 5.35, "{{primary_workflow}}", 1.8),
        (2.67, 5.35, "DL verification", 1.35),
        (4.18, 5.35, "RC checks", 1.05),
        (0.72, 5.78, "Aadhaar + PAN", 1.35),
        (2.33, 5.78, "For {{prospect_role}}", 1.7),
    ]:
        add_chip(slide, x, y, text, width)
    add_panel(slide, 7.0, 1.02, 3.95, 4.72, fill="panel_soft")
    add_text(slide, 7.27, 1.26, 3.25, 0.58, "What this changes for {{lead_company_name}}", size=14.5, bold=True, color="ink")
    add_metric(slide, 7.27, 2.02, 1.55, 1.32, "Masked KYC", "Aadhaar data is extracted while sensitive digits remain protected")
    add_metric(slide, 9.05, 2.02, 1.55, 1.32, "API ready", "DL and RC records can return Parivahan-backed status signals")
    add_metric(slide, 7.27, 3.5, 1.55, 1.32, "Ops ready", "JSON payloads, confidence flags, and review status")
    add_metric(slide, 9.05, 3.5, 1.55, 1.32, "Faster activation", "Clean drivers and vehicles move to approval faster")
    add_band(slide, 7.05, 5.92, 3.9, 1.12, "Positioning", "Replace manual onboarding review with one verification layer for {{lead_company_name}}.")


def build_slide_2():
    slide = new_slide("Why Mobility Ops Break Down", 3)
    add_eyebrow(slide, 0.72, 1.0, "THE OPERATIONAL PROBLEM", w=2.35)
    add_text(slide, 0.72, 1.48, 6.0, 0.98, "Driver and vehicle onboarding is still document-heavy, slow, and error-prone.", size=21.5, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.58, 5.9, 0.48, "For {{lead_company_name}}, {{pain_point_1}}, {{pain_point_2}}, and {{pain_point_3}} can slow down {{primary_workflow}}.", size=10.9, color="ink_soft")
    cards = [
        ("KYC", "Identity capture fails in the field", "Blurred Aadhaar uploads slow approvals and compliance checks."),
        ("Vehicle Ops", "DL and RC data is retyped manually", "Manual keying creates errors in onboarding records."),
        ("Risk", "Validation happens too late", "Parivahan checks often sit outside the core onboarding flow."),
        ("Scale", "Reviewer queues pile up fast", "Launches and hiring pushes create instant backlogs."),
    ]
    coords = [(0.72, 3.28), (3.55, 3.28), (0.72, 5.1), (3.55, 5.1)]
    for (tag, title, body), (x, y) in zip(cards, coords):
        add_card(slide, x, y, 2.55, 1.58, tag, title, body, tag_w=1.0)
    add_panel(slide, 6.85, 1.52, 3.95, 2.2)
    add_text(slide, 7.08, 1.78, 3.42, 0.33, "Typical failure pattern", size=17, bold=True, color="ink")
    add_multiline(slide, 7.1, 2.24, 3.35, 1.05, [
        "Mixed-quality Aadhaar, PAN, DL, and RC uploads",
        "Raw OCR without normalization or validation",
        "Manual cleanup across multiple tabs",
        "Slow approvals and lower driver conversion",
    ], size=10.3, bullet=True)
    add_panel(slide, 6.85, 4.1, 3.95, 2.12, fill="panel_soft")
    add_text(slide, 7.08, 4.35, 3.42, 0.33, "What operators need instead", size=17, bold=True, color="ink")
    add_multiline(slide, 7.1, 4.8, 3.35, 1.0, [
        "One layer for extraction, masking, and validation",
        "Parivahan-backed DL and RC checks",
        "Review only low-confidence cases",
        "Audit-ready outputs for compliance teams",
    ], size=10.3, bullet=True)


def add_signal_block(slide, x, y, w, h, label, value):
    add_panel(slide, x, y, w, h, fill="panel")
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.22, label, size=8.2, bold=True, color="accent_dark")
    add_text(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.58, value, size=9.2, color="ink_soft")


def add_field_row(slide, x, y, label, value, w=2.9):
    add_text(slide, x, y, 0.74, 0.18, label.upper(), size=5.4, bold=True, color="ink_soft")
    add_text(slide, x + 0.78, y - 0.01, w - 0.78, 0.2, value, size=7.5, bold=True, color="ink")


def add_sample_document(slide, x, y, w, h, tag, title, subtitle, rows, fill="panel"):
    add_panel(slide, x, y, w, h, fill=fill)
    add_tag(slide, x + 0.16, y + 0.14, tag, w=0.72)
    add_text(slide, x + 0.98, y + 0.15, w - 1.14, 0.22, title, size=11.4, bold=True, color="ink")
    add_text(slide, x + 0.16, y + 0.49, w - 0.32, 0.2, subtitle, size=7.2, color="ink_soft")
    cy = y + 0.82
    for label, value in rows:
        add_field_row(slide, x + 0.2, cy, label, value, w=w - 0.4)
        cy += 0.3


def add_compact_json(slide, x, y, w, h, title, lines):
    add_panel(slide, x, y, w, h, fill="navy", line="navy")
    add_text(slide, x + 0.18, y + 0.15, w - 0.36, 0.2, title, size=9.8, bold=True, color="white")
    add_multiline(slide, x + 0.18, y + 0.45, w - 0.36, h - 0.54, lines, size=7.0, color="white", bullet=False)


def build_slide_4_visual_proof():
    slide = new_slide("Visual Proof", 4)
    add_eyebrow(slide, 0.72, 1.0, "DOCUMENTS + API RESPONSES", w=2.9)
    add_text(slide, 0.72, 1.46, 8.75, 0.62, "Visual proof: KYC images, extracted data, and Parivahan status.", size=19.8, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.17, 8.55, 0.3, "Masked Aadhaar and PAN examples flow into structured extraction, while DL and RC details return Parivahan-backed verification responses.", size=9.8, color="ink_soft")

    add_panel(slide, 0.72, 2.75, 4.95, 4.22)
    add_text(slide, 0.95, 2.98, 3.55, 0.25, "Sample uploaded documents", size=13.8, bold=True, color="ink")
    add_sample_document(slide, 0.95, 3.45, 2.15, 2.15, "KYC", "Aadhaar", "Masked identity image", [
        ("Name", "RAVI KUMAR"),
        ("Aadhaar", "XXXX XXXX 4321"),
        ("DOB", "19XX"),
        ("Address", "Bengaluru"),
    ], fill="panel_soft")
    add_sample_document(slide, 3.25, 3.45, 2.15, 2.15, "PAN", "PAN Card", "Tax ID image", [
        ("Name", "RAVI KUMAR"),
        ("PAN", "ABCDE1234F"),
        ("Status", "Name match"),
        ("DOB", "19XX"),
    ])
    add_panel(slide, 0.95, 5.92, 4.45, 0.62, fill="panel_soft")
    add_text(slide, 1.14, 6.08, 1.5, 0.18, "Parivahan inputs", size=8.6, bold=True, color="accent_dark")
    add_text(slide, 2.56, 6.08, 2.55, 0.18, "DL: KA05 20XX 1234567 | RC: KA-05-AB-1234", size=7.4, color="ink_soft")

    add_panel(slide, 5.95, 2.75, 2.25, 2.05, fill="panel_soft")
    add_text(slide, 6.16, 2.98, 1.78, 0.25, "Extracted data", size=13.2, bold=True, color="ink")
    add_multiline(slide, 6.18, 3.38, 1.75, 1.05, [
        "driver_name: Ravi Kumar",
        "aadhaar: XXXX-XXXX-4321",
        "pan: ABCDE1234F",
        "dl_no: KA05 20XX 1234567",
        "rc_no: KA-05-AB-1234",
    ], size=7.4, bullet=False)

    add_panel(slide, 8.42, 2.75, 2.25, 2.05)
    add_text(slide, 8.63, 2.98, 1.72, 0.25, "Fracto signals", size=13.2, bold=True, color="ink")
    add_multiline(slide, 8.65, 3.38, 1.72, 1.05, [
        "aadhaar_masked: true",
        "pan_name_match: true",
        "field_confidence: high",
        "review_status: ready",
        "route_to: {{downstream_system}}",
    ], size=7.4, bullet=False)

    add_compact_json(slide, 5.95, 5.1, 2.25, 1.55, "Parivahan DL API", [
        "status: active",
        "holder_match: true",
        "vehicle_class: LMV/TR",
        "valid_till: 2028-12",
    ])
    add_compact_json(slide, 8.42, 5.1, 2.25, 1.55, "Parivahan RC API", [
        "status: verified",
        "owner_match: true",
        "fitness: active",
        "insurance: valid",
    ])


def build_slide_4():
    slide = new_slide("Workflow Fit", 5)
    add_eyebrow(slide, 0.72, 1.0, "END-TO-END FLOW", w=1.95)
    add_text(slide, 0.72, 1.48, 8.4, 0.72, "Fracto can sit inside the {{lead_company_name}} driver onboarding workflow.", size=21.5, bold=True, color="ink", font="Space Grotesk")
    steps = [
        ("01", "Capture", "Drivers upload Aadhaar, PAN, DL, RC, and required KYC files."),
        ("02", "Parse and Mask", "Fracto classifies files, extracts fields, and masks Aadhaar before downstream use."),
        ("03", "Validate", "DL and RC checks run through Parivahan-backed validation and business rules."),
        ("04", "Decide", "Clean records move to {{downstream_system}} while exceptions route to {{review_team}}."),
    ]
    for idx, (no, title, body) in enumerate(steps):
        x = 0.72 + idx * 2.62
        add_panel(slide, x, 2.42, 2.35, 1.92)
        add_text(slide, x + 0.18, 2.64, 0.55, 0.35, no, size=19, bold=True, color="accent", font="Space Grotesk")
        add_text(slide, x + 0.18, 3.05, 1.9, 0.3, title, size=13.5, bold=True, color="ink")
        add_text(slide, x + 0.18, 3.43, 1.95, 0.58, body, size=9.4, color="ink_soft")
        if idx < 3:
            add_text(slide, x + 2.37, 3.12, 0.22, 0.25, ">", size=15, bold=True, color="accent", align=PP_ALIGN.CENTER)
    add_panel(slide, 0.72, 5.0, 5.1, 1.78)
    add_text(slide, 0.95, 5.22, 4.4, 0.28, "High-value insertion points", size=15.5, bold=True, color="ink")
    add_multiline(slide, 0.98, 5.62, 4.45, 0.98, [
        "Driver onboarding", "Vehicle onboarding", "Compliance refreshes", "Partner onboarding", "Central review queues"
    ], size=8.7, bullet=True)
    add_panel(slide, 6.2, 5.0, 4.55, 1.78, fill="panel_soft")
    add_text(slide, 6.43, 5.22, 3.8, 0.28, "Decision signals", size=15.5, bold=True, color="ink")
    add_multiline(slide, 6.46, 5.62, 3.8, 0.98, [
        "Aadhaar masked before sharing", "PAN name match or mismatch", "DL and RC Parivahan responses", "Field confidence and exception flags"
    ], size=8.7, bullet=True)


def build_slide_6():
    slide = new_slide("Operator Impact", 6)
    add_eyebrow(slide, 0.72, 1.0, "BUSINESS OUTCOMES", w=2.1)
    add_text(slide, 0.72, 1.48, 8.55, 1.28, "The value case for {{prospect_role}} at {{lead_company_name}}.", size=19.8, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.88, 8.35, 0.32, "Help improve {{outcome_1}}, {{outcome_2}}, and {{outcome_3}} across {{primary_workflow}}.", size=10.2, color="ink_soft")
    impact = [
        ("Faster driver activation", "Clean records move from upload to approval faster."),
        ("Lower manual review", "Reviewers focus only on low-confidence or mismatched records."),
        ("Stronger compliance", "Aadhaar masking, Parivahan checks, and audit-ready payloads in one layer."),
    ]
    for idx, (title, body) in enumerate(impact):
        x = 0.72 + idx * 3.42
        add_panel(slide, x, 3.32, 3.05, 1.23)
        add_text(slide, x + 0.18, 3.52, 2.55, 0.3, title, size=14.5, bold=True, color="ink", font="Space Grotesk")
        add_text(slide, x + 0.18, 3.94, 2.55, 0.33, body, size=8.4, color="ink_soft")
    add_panel(slide, 0.72, 4.9, 4.95, 1.55)
    add_text(slide, 0.95, 5.12, 4.3, 0.28, "Common workflow bottlenecks", size=15.5, bold=True, color="ink")
    add_multiline(slide, 0.98, 5.5, 4.15, 0.72, [
        "Aadhaar masking is handled manually", "PAN, DL, and RC fields are retyped", "Parivahan checks happen in separate tabs", "Review queues grow during onboarding spikes"
    ], size=9.7, bullet=True)
    add_panel(slide, 6.05, 4.9, 4.7, 1.55, fill="panel_soft")
    add_text(slide, 6.28, 5.12, 4.05, 0.28, "What Fracto enables", size=15.5, bold=True, color="ink")
    add_multiline(slide, 6.31, 5.5, 3.9, 0.72, [
        "Structured KYC and vehicle outputs", "Source-system validation signals", "Reviewer queues only for exceptions", "Faster activation with cleaner audit trail"
    ], size=9.7, bullet=True)


def build_slide_8():
    slide = new_slide("Book Demo", 7)
    add_eyebrow(slide, 0.72, 1.0, "NEXT STEP", w=1.3)
    add_text(slide, 0.72, 1.48, 5.4, 1.55, "{{prospect_name}}, want to review this on a few {{lead_company_name}} onboarding files?", size=20.5, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 3.15, 5.15, 0.48, "Book a demo to review Aadhaar masking, PAN extraction, DL/RC extraction, Parivahan checks, and exception routing on sample onboarding documents.", size=10.6, color="ink_soft")
    add_panel(slide, 0.72, 4.02, 5.2, 1.68)
    add_text(slide, 0.95, 4.24, 3.8, 0.28, "Suggested demo scope", size=14, bold=True, color="ink")
    add_multiline(slide, 0.98, 4.65, 4.45, 0.76, [
        "Input bundle: {{sample_document_types}}",
        "Flow: {{demo_flow}}",
        "Success metric: {{success_metric}}",
    ], size=9.2, bullet=True)
    add_band(slide, 0.72, 5.9, 5.2, 1.05, "Demo Focus", "Review real onboarding files first. Then scope the right integration path.")
    add_panel(slide, 6.45, 1.28, 4.3, 5.55, fill="navy", line="navy")
    add_eyebrow(slide, 6.75, 1.68, "TALK TO FRACTO", w=1.65)
    add_text(slide, 6.75, 2.25, 3.55, 0.65, "Book a demo and review capabilities.", size=22, bold=True, color="white", font="Space Grotesk")
    contact = [("Email", "sales@fracto.tech"), ("Phone", "+91 84483 33277"), ("Website", "fracto.tech")]
    y = 3.22
    for label, value in contact:
        add_panel(slide, 6.78, y, 3.65, 0.58, fill="accent_dark", line="accent_dark")
        add_text(slide, 6.95, y + 0.08, 0.9, 0.18, label.upper(), size=7.5, bold=True, color="white")
        add_text(slide, 7.82, y + 0.06, 2.35, 0.22, value, size=12, bold=True, color="white", font="Space Grotesk")
        y += 0.78
    add_text(slide, 6.78, 5.92, 3.6, 0.26, "Next step: review a small {{primary_workflow}} sample set.", size=8.8, color="white")


def build_slide_5_references():
    slide = new_slide("Reference Set", 2)
    add_eyebrow(slide, 0.72, 1.0, "MOBILITY OPERATIONS", w=2.05)
    add_text(slide, 0.72, 1.48, 5.8, 1.04, "Fracto supports document-heavy mobility and fleet workflows.", size=21.0, bold=True, color="ink", font="Space Grotesk")
    add_text(slide, 0.72, 2.58, 5.65, 0.4, "For {{lead_company_name}}, the same layer can support driver onboarding, vehicle checks, KYC masking, and recurring compliance review.", size=10.5, color="ink_soft")

    logos = [
        ("clients/routematic.png", "Routematic"),
        ("clients/goamiles-tight.png", "GoaMiles"),
        ("clients/alt-mobility.png", "Alt Mobility"),
        ("clients/whistledrive.png", "WhistleDrive"),
        ("clients/turno.png", "Turno"),
        ("clients/gozo-cabs.png", "Gozo Cabs"),
    ]
    add_panel(slide, 0.72, 3.35, 5.75, 2.55, fill="panel_soft")
    add_text(slide, 0.95, 3.56, 3.8, 0.25, "Relevant mobility and fleet workflows", size=13.4, bold=True, color="ink")
    for idx, (path, name) in enumerate(logos):
        x = 0.95 + (idx % 3) * 1.78
        y = 4.02 + (idx // 3) * 0.9
        add_logo_card(slide, x, y, 1.5, 0.64, path, name)

    add_panel(slide, 6.85, 1.52, 3.95, 2.22)
    add_text(slide, 7.08, 1.76, 3.42, 0.68, "Why this matters for {{lead_company_name}}", size=15.2, bold=True, color="ink")
    add_multiline(slide, 7.1, 2.58, 3.35, 0.86, [
        "High onboarding document volume",
        "Sensitive KYC data needs masking",
        "Driver and vehicle records need validation",
        "Faster activation improves supply availability",
    ], size=8.9, bullet=True)
    add_panel(slide, 6.85, 4.08, 3.95, 2.02, fill="panel_soft")
    add_text(slide, 7.08, 4.32, 3.42, 0.3, "What your team can expect", size=15.8, bold=True, color="ink")
    add_multiline(slide, 7.1, 4.74, 3.35, 0.94, [
        "Cleaner driver and vehicle records",
        "Lower manual review effort",
        "One API layer for onboarding files",
        "Audit-ready outputs for operations",
    ], size=9.8, bullet=True)


for builder in [
    build_slide_1,
    build_slide_5_references,
    build_slide_2,
    build_slide_4_visual_proof,
    build_slide_4,
    build_slide_6,
    build_slide_8,
]:
    builder()

props = prs.core_properties
props.title = "Fracto Mobility Driver Onboarding Dynamic Sales Deck"
props.subject = "Personalized client-facing sales deck with editable variables"
props.author = "Fracto"
TMP_OUT = OUT.with_suffix(".tmp.pptx")
if TMP_OUT.exists():
    TMP_OUT.unlink()
prs.save(TMP_OUT)
TMP_OUT.replace(OUT)
print(OUT)
